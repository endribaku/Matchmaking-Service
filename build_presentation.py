"""
Build the project presentation as a .pptx file.

  Run:    python3 build_presentation.py
  Output: Matchmaking_Service_Presentation.pptx

Dark/orange theme to match the app. 16:9 widescreen.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- theme ----------------------------------------------------------------

BG     = RGBColor(0x0d, 0x0f, 0x12)
BG2    = RGBColor(0x16, 0x1a, 0x20)
BG3    = RGBColor(0x1d, 0x23, 0x2b)
FG     = RGBColor(0xe8, 0xea, 0xed)
MUTED  = RGBColor(0x8a, 0x8f, 0x97)
ACCENT = RGBColor(0xf5, 0x9e, 0x0b)
GOOD   = RGBColor(0x10, 0xb9, 0x81)
BORDER = RGBColor(0x2a, 0x30, 0x3a)

SANS = "Calibri"
MONO = "Consolas"

# ---- presentation setup ---------------------------------------------------

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ---- helpers --------------------------------------------------------------

def add_slide(page_num=None, total=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # top accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.12))
    strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    strip.shadow.inherit = False
    # page number bottom-right
    if page_num is not None and total is not None:
        text(s, Inches(12.4), Inches(7.1), Inches(0.85), Inches(0.3),
             f"{page_num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    return s

def text(slide, left, top, width, height, content, *,
         size=18, bold=False, color=FG, align=PP_ALIGN.LEFT,
         font=SANS, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """content can be a string or list of (text, opts_dict) tuples for mixed runs.
       If list of strings, each is a separate paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(content, str):
        paragraphs = [content]
    else:
        paragraphs = content

    for i, p_content in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = p_content
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb

def title_block(slide, kicker, title):
    text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
         kicker, size=11, color=ACCENT, bold=True)
    text(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
         title, size=34, color=FG, bold=True)
    # underline accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.55), Inches(0.6), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

def card(slide, left, top, width, height, *, fill=BG2, border=BORDER):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(0.75)
    c.shadow.inherit = False
    # tighten corner radius
    c.adjustments[0] = 0.05
    return c

# ---- slides ---------------------------------------------------------------

TOTAL = 11  # filled in after we know the count; kept manual for clarity

def slide_title():
    s = add_slide()  # no page number on cover
    # accent dot in top-left
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(2.4),
                             Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background(); dot.shadow.inherit = False

    text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(0.5),
         "PARALLEL PROGRAMMING PROJECT",
         size=14, color=ACCENT, bold=True)
    text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(1.4),
         "Matchmaking Service",
         size=64, color=FG, bold=True)
    text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.7),
         "A 5v5 matchmaking system built with java.util.concurrent",
         size=22, color=MUTED)
    # footer
    text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Student: ____________________     |     Course: Parallel Programming     |     Epoka University, 2026",
         size=13, color=MUTED)

def slide_what():
    s = add_slide(2, TOTAL)
    title_block(s, "OVERVIEW", "What is it?")

    text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
         "A small 5v5 matchmaking service inspired by competitive games.",
         size=22, color=FG, bold=True)

    bullets = [
        "Players join a queue with an MMR (skill rating).",
        "A background matchmaker continuously pairs 10 skill-close players into two balanced teams.",
        "Bots are constantly added to the queue, so the service is always alive.",
        "Same kind of problem solved by CS2, Valorant, LoL and Dota 2 — at student scale.",
    ]
    for i, b in enumerate(bullets):
        # bullet dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.65), Inches(3.0 + i*0.85), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background(); dot.shadow.inherit = False
        text(s, Inches(1.0), Inches(2.93 + i*0.85), Inches(11.5), Inches(0.6),
             b, size=18, color=FG)

def slide_stack():
    s = add_slide(3, TOTAL)
    title_block(s, "TECHNOLOGY", "Tech stack")

    # two cards side by side
    cw = Inches(5.85); ch = Inches(4.0)
    y  = Inches(2.2)

    # Backend card
    card(s, Inches(0.6), y, cw, ch)
    text(s, Inches(0.9), y + Inches(0.3), cw - Inches(0.4), Inches(0.4),
         "BACKEND", size=13, color=ACCENT, bold=True)
    text(s, Inches(0.9), y + Inches(0.7), cw - Inches(0.4), Inches(0.6),
         "Java 17+", size=28, color=FG, bold=True)
    backend_lines = [
        "• com.sun.net.httpserver.HttpServer  (built-in JDK)",
        "• java.util.concurrent  (all parallelism)",
        "• Jackson  (JSON serialization)",
        "• Maven  (build)",
    ]
    text(s, Inches(0.9), y + Inches(1.5), cw - Inches(0.4), Inches(2.4),
         backend_lines, size=15, color=FG, line_spacing=1.4)

    # Frontend card
    fx = Inches(6.85)
    card(s, fx, y, cw, ch)
    text(s, fx + Inches(0.3), y + Inches(0.3), cw - Inches(0.4), Inches(0.4),
         "FRONTEND", size=13, color=ACCENT, bold=True)
    text(s, fx + Inches(0.3), y + Inches(0.7), cw - Inches(0.4), Inches(0.6),
         "React 18 + Vite", size=28, color=FG, bold=True)
    frontend_lines = [
        "• Plain JavaScript  (no TypeScript)",
        "• 4 components  (Login, Lobby, Queue, Match)",
        "• LiveStats sidebar  (always visible)",
        "• Polls /api/stats every 1 s",
    ]
    text(s, fx + Inches(0.3), y + Inches(1.5), cw - Inches(0.4), Inches(2.4),
         frontend_lines, size=15, color=FG, line_spacing=1.4)

    text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "No frameworks. No databases. Just the JDK + React. Everything lives in memory while the server runs.",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)

def slide_components():
    s = add_slide(4, TOTAL)
    title_block(s, "ARCHITECTURE", "Four parallel components")

    components = [
        ("1", "HTTP request pool",     "8 worker threads",
         "Serves the REST API. Multiple browsers can register, queue and poll at the same time."),
        ("2", "Matchmaker tick",       "1 thread, every 500 ms",
         "Drains the queue, sorts by MMR, finds 10 skill-close players, builds two balanced teams."),
        ("3", "Bot populator",         "1 thread, every 500 ms",
         "Keeps the queue full of bots so the service is always alive. Target scales with worker count."),
        ("4", "Match-setup pool",      "N workers (configurable)",
         "After a match is formed, setup work runs here. More workers ⇒ more matches prepared in parallel."),
    ]
    cw = Inches(6.05); ch = Inches(2.25); gx = Inches(0.6); gy = Inches(1.95)
    for i, (num, name, sub, desc) in enumerate(components):
        col, row = i % 2, i // 2
        x = gx + col * (cw + Inches(0.15))
        y = gy + row * (ch + Inches(0.2))
        card(s, x, y, cw, ch)
        # number circle
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.3),
                                Inches(0.55), Inches(0.55))
        nc.fill.solid(); nc.fill.fore_color.rgb = ACCENT
        nc.line.fill.background(); nc.shadow.inherit = False
        text(s, x + Inches(0.25), y + Inches(0.32), Inches(0.55), Inches(0.5),
             num, size=18, color=BG, bold=True, align=PP_ALIGN.CENTER)
        # name + sub + desc
        text(s, x + Inches(1.0), y + Inches(0.3), cw - Inches(1.2), Inches(0.5),
             name, size=20, color=FG, bold=True)
        text(s, x + Inches(1.0), y + Inches(0.78), cw - Inches(1.2), Inches(0.35),
             sub, size=12, color=ACCENT)
        text(s, x + Inches(0.4), y + Inches(1.25), cw - Inches(0.6), Inches(0.95),
             desc, size=14, color=FG, line_spacing=1.3)

def slide_pipeline():
    s = add_slide(5, TOTAL)
    title_block(s, "PIPELINE", "How a player flows through the thread pools")

    text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Four threads, two shared structures. Every horizontal arrow crosses a thread boundary.",
         size=15, color=MUTED)

    # 4 stage layout: producers → queue → matchmaker → setup workers
    lane_w  = Inches(2.75)
    arrow_w = Inches(0.45)
    margin  = Inches(0.5)
    gap     = Inches(0.05)

    y_head = Inches(2.65)
    y_box  = Inches(3.05)
    box_h  = Inches(2.7)

    def lane_x(i):
        return margin + i * (lane_w + arrow_w + 2 * gap)

    def stage(i, kicker, kicker_color, title, body_lines, fill=BG2):
        x = lane_x(i)
        # kicker
        text(s, x, y_head, lane_w, Inches(0.32),
             kicker, size=11, color=kicker_color, bold=True, align=PP_ALIGN.CENTER)
        # box
        card(s, x, y_box, lane_w, box_h, fill=fill)
        # title inside box
        text(s, x + Inches(0.15), y_box + Inches(0.2), lane_w - Inches(0.3), Inches(0.45),
             title, size=15, color=FG, bold=True, align=PP_ALIGN.CENTER)
        # divider
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            x + Inches(0.7), y_box + Inches(0.78), lane_w - Inches(1.4), Inches(0.02))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background(); bar.shadow.inherit = False
        # body lines
        text(s, x + Inches(0.2), y_box + Inches(1.0), lane_w - Inches(0.4),
             box_h - Inches(1.1),
             body_lines, size=12, color=FG, line_spacing=1.45)

    def arrow(i):
        # arrow between stage i and stage i+1
        x = lane_x(i) + lane_w + gap
        y = y_box + box_h / 2 - Inches(0.2)
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, arrow_w, Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = ACCENT
        a.line.fill.background(); a.shadow.inherit = False

    stage(0, "1 — PRODUCERS", ACCENT, "HTTP join + populator", [
        "• 8 HTTP workers push tickets",
        "  when players click Find Match",
        "• Bot populator pushes bots",
        "  every 500 ms to keep queue full",
    ])
    arrow(0)
    stage(1, "2 — SHARED QUEUE", GOOD, "LinkedBlockingQueue", [
        "• Thread-safe FIFO",
        "• Producers call offer()",
        "• Consumer (matchmaker)",
        "  calls poll() once per tick",
    ], fill=BG3)
    arrow(1)
    stage(2, "3 — MATCHMAKER TICK", ACCENT, "1 thread · every 500 ms", [
        "• Drain queue → sorted pool",
        "• Sort by MMR ascending",
        "• Sliding window of 10",
        "• Snake-draft team balance",
        "• matches.put + setupPool.submit",
    ])
    arrow(2)
    stage(3, "4 — SETUP POOL", ACCENT, "N workers · resizable live", [
        "• Each worker prepares one match",
        "  (simulated 800 ms allocation)",
        "• With pool size N, up to N matches",
        "  set up at the same time",
        "• Dropdown: 1 / 2 / 4 / 8",
    ])

    # Bottom callout: status read path (lock-free)
    y_status = Inches(6.15)
    text(s, Inches(0.6), y_status, Inches(12), Inches(0.32),
         "PARALLEL READ PATH", size=11, color=ACCENT, bold=True)
    card(s, Inches(0.6), y_status + Inches(0.35), Inches(12.1), Inches(0.7), fill=BG2)
    text(s, Inches(0.9), y_status + Inches(0.5), Inches(11.5), Inches(0.45),
         "GET /api/queue/status  →  ConcurrentHashMap.get(playerId)  →  match payload  "
         "(no lock; all polls in parallel)",
         size=14, color=FG, font=MONO, align=PP_ALIGN.CENTER)

def slide_primitives():
    s = add_slide(6, TOTAL)
    title_block(s, "CONCURRENCY", "java.util.concurrent primitives used")

    rows = [
        ("LinkedBlockingQueue<QueueTicket>", "Producer/consumer hand-off — HTTP threads push, matchmaker drains."),
        ("ScheduledExecutorService × 2",     "Matchmaker tick and bot populator on separate threads."),
        ("ThreadPoolExecutor × 2",           "HTTP request pool and match-setup pool (resizable live)."),
        ("ConcurrentHashMap",                "Lock-free reads of player & match registries under heavy polling."),
        ("ConcurrentLinkedQueue / AtomicLong", "Lock-free counters: tick count, match IDs, bots spawned."),
        ("ReentrantLock",                    "Critical section around the sorted waiting pool inside one tick."),
        ("volatile + daemon threads",        "Cross-thread visibility for last-tick duration; clean shutdown."),
    ]
    y0 = Inches(2.0)
    rh = Inches(0.62)
    for i, (name, desc) in enumerate(rows):
        y = y0 + i * rh
        # alternating row background
        bg = BG2 if i % 2 == 0 else BG3
        row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.6), y, Inches(12.1), rh - Inches(0.04))
        row.fill.solid(); row.fill.fore_color.rgb = bg
        row.line.fill.background(); row.shadow.inherit = False
        text(s, Inches(0.85), y + Inches(0.12), Inches(4.6), Inches(0.45),
             name, size=14, color=ACCENT, bold=True, font=MONO)
        text(s, Inches(5.6),  y + Inches(0.13), Inches(7.0), Inches(0.45),
             desc, size=14, color=FG)

def slide_algorithm():
    s = add_slide(7, TOTAL)
    title_block(s, "ALGORITHM", "How the matchmaker forms a match")

    text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Every 500 ms, on a dedicated background thread:",
         size=16, color=MUTED)

    steps = [
        ("DRAIN",    "Pull new arrivals from the BlockingQueue into a sorted pool."),
        ("FILTER",   "Drop tickets for players who cancelled."),
        ("SORT",     "Sort the pool by MMR ascending — adjacent players are skill-close."),
        ("SCAN",     "Slide a window of 10. If the MMR spread fits the tolerance, form a match."),
        ("RELAX",    "Tolerance starts at 300 MMR and widens 100 MMR / second waited (cap 1500)."),
        ("BALANCE",  "Sort the 10 by MMR desc, assign A,B,B,A,A,B,B,A,A,B — snake draft."),
        ("HAND OFF", "Submit setupMatch() to the parallel setup pool. Tick continues without blocking."),
    ]
    y0 = Inches(2.7)
    for i, (kw, desc) in enumerate(steps):
        y = y0 + i * Inches(0.6)
        text(s, Inches(0.75), y, Inches(0.6), Inches(0.4),
             f"{i+1}.", size=18, color=ACCENT, bold=True)
        text(s, Inches(1.25), y - Inches(0.02), Inches(1.6), Inches(0.45),
             kw, size=14, color=ACCENT, bold=True, font=MONO)
        text(s, Inches(3.0), y - Inches(0.02), Inches(9.5), Inches(0.5),
             desc, size=15, color=FG)

def slide_balance():
    s = add_slide(8, TOTAL)
    title_block(s, "TEAM BALANCE", "Snake draft keeps teams fair")

    text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Sort 10 players by MMR (highest first), then assign in pattern A B B A A B B A A B.",
         size=15, color=MUTED)

    # Two team cards
    cw = Inches(5.7); ch = Inches(3.7); y = Inches(2.65)

    def team_card(x, label, players, total, color_accent):
        card(s, x, y, cw, ch)
        text(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(0.4),
             label, size=13, color=color_accent, bold=True)
        text(s, x + Inches(0.3), y + Inches(0.6), cw - Inches(0.6), Inches(0.5),
             f"{total} MMR", size=22, color=FG, bold=True)
        for i, (name, mmr) in enumerate(players):
            ry = y + Inches(1.3 + i * 0.4)
            text(s, x + Inches(0.4), ry, Inches(3.0), Inches(0.35),
                 name, size=14, color=FG)
            text(s, x + cw - Inches(1.3), ry, Inches(1.0), Inches(0.35),
                 f"{mmr} MMR", size=14, color=MUTED, align=PP_ALIGN.RIGHT)

    team_a = [("Player #1",  2000),
              ("Player #4",  1800),
              ("Player #5",  1700),
              ("Player #8",  1500),
              ("Player #9",  1450)]
    team_b = [("Player #2",  1900),
              ("Player #3",  1850),
              ("Player #6",  1650),
              ("Player #7",  1600),
              ("Player #10", 1400)]
    team_card(Inches(0.6), "TEAM A",  team_a, sum(m for _, m in team_a), ACCENT)
    team_card(Inches(7.05),"TEAM B",  team_b, sum(m for _, m in team_b), GOOD)

    text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         f"Total difference: {abs(sum(m for _, m in team_a) - sum(m for _, m in team_b))} MMR.  "
         "Typical match: under 100 MMR across two 5-stacks.",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)

def slide_ui():
    s = add_slide(9, TOTAL)
    title_block(s, "INTERFACE", "Live parallelism, visible in the UI")

    text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "An always-visible sidebar shows the four threads working in real time.",
         size=17, color=FG)

    # Left: feature list
    features = [
        ("Two big counters",   "Players in queue and total matches formed."),
        ("Pulse indicators",   "Green dot lights up when a thread is busy."),
        ("Capacity bars",      "Visualize how many setup workers are running right now."),
        ("Live dropdown",      "Resize the setup pool live: 1, 2, 4 or 8 workers."),
        ("Visible speedup",    "Switch from 1 to 8 workers — match throughput jumps ~6×."),
    ]
    y0 = Inches(2.7)
    for i, (k, v) in enumerate(features):
        y = y0 + i * Inches(0.75)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.65), y + Inches(0.12), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background(); dot.shadow.inherit = False
        text(s, Inches(1.05), y, Inches(5.4), Inches(0.4),
             k, size=15, color=ACCENT, bold=True)
        text(s, Inches(1.05), y + Inches(0.35), Inches(5.6), Inches(0.4),
             v, size=14, color=FG)

    # Right: mock sidebar card
    x = Inches(7.5); y = Inches(2.4); w = Inches(5.2); h = Inches(4.7)
    card(s, x, y, w, h)
    text(s, x + Inches(0.35), y + Inches(0.3), w, Inches(0.3),
         "LIVE SYSTEM ACTIVITY", size=11, color=ACCENT, bold=True)
    # big numbers
    bx = x + Inches(0.35); by = y + Inches(0.75)
    bw = (w - Inches(0.9)) / 2
    bh = Inches(0.95)
    for i, (n, label) in enumerate([("78", "in queue"), ("42", "matches formed")]):
        rx = bx + i * (bw + Inches(0.2))
        card(s, rx, by, bw, bh, fill=BG3, border=BORDER)
        text(s, rx, by + Inches(0.05), bw, Inches(0.55),
             n, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        text(s, rx, by + Inches(0.55), bw, Inches(0.4),
             label, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # thread rows
    rows = [
        ("Matchmaker tick",   "tick #142 · 3 ms",          True),
        ("Bot populator",     "264 bots spawned",          True),
        ("Match-setup pool",  "4 / 4 workers busy",        True),
        ("HTTP request pool", "1 / 8 workers busy",        True),
    ]
    ry = by + bh + Inches(0.4)
    for name, sub, active in rows:
        # dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
            bx + Inches(0.05), ry + Inches(0.13), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOOD if active else BORDER
        dot.line.fill.background(); dot.shadow.inherit = False
        text(s, bx + Inches(0.35), ry, w - Inches(1.0), Inches(0.3),
             name, size=12, color=FG, bold=True)
        text(s, bx + Inches(0.35), ry + Inches(0.25), w - Inches(1.0), Inches(0.3),
             sub, size=10, color=MUTED, font=MONO)
        ry += Inches(0.55)

def slide_run():
    s = add_slide(10, TOTAL)
    title_block(s, "RUN", "How to run the demo")

    # Step 1
    text(s, Inches(0.6), Inches(2.05), Inches(12), Inches(0.5),
         "1.  Start the backend  (in one terminal)",
         size=18, color=ACCENT, bold=True)
    card(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.8), fill=BG2)
    text(s, Inches(0.9), Inches(2.78), Inches(11.4), Inches(0.45),
         "cd backend && mvn compile exec:java",
         size=18, color=FG, font=MONO)

    # Step 2
    text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5),
         "2.  Start the frontend  (in another terminal)",
         size=18, color=ACCENT, bold=True)
    card(s, Inches(0.6), Inches(4.25), Inches(12), Inches(0.8), fill=BG2)
    text(s, Inches(0.9), Inches(4.43), Inches(11.4), Inches(0.45),
         "cd frontend && npm install && npm run dev",
         size=18, color=FG, font=MONO)

    # Step 3
    text(s, Inches(0.6), Inches(5.35), Inches(12), Inches(0.5),
         "3.  Open the app and play",
         size=18, color=ACCENT, bold=True)
    card(s, Inches(0.6), Inches(5.9), Inches(12), Inches(1.1), fill=BG2)
    bullets = [
        "Open  http://localhost:5173  — enter any name to register.",
        "Click 'Find Match' — you'll be matched within a second (the queue is full of bots).",
        "Try the 'Match-setup workers' dropdown: switch from 1 to 8 and watch the capacity bar fill up.",
    ]
    for i, b in enumerate(bullets):
        text(s, Inches(0.95), Inches(5.95) + i * Inches(0.32), Inches(11.4), Inches(0.32),
             "•  " + b, size=12, color=FG)

def slide_summary():
    s = add_slide(11, TOTAL)
    title_block(s, "SUMMARY", "What this project demonstrates")

    points = [
        ("Real-world domain",
         "Game matchmaking — the same class of problem CS2, Valorant and LoL solve."),
        ("Four cooperating thread pools",
         "HTTP workers, matchmaker tick, bot populator, match-setup pool."),
        ("Seven concurrency primitives",
         "BlockingQueue, ScheduledExecutorService, ThreadPoolExecutor, ConcurrentHashMap, AtomicLong, ReentrantLock, volatile."),
        ("Live visualization",
         "The UI shows every thread pool busy in real time."),
        ("Tunable parallelism",
         "Resize the setup pool from 1 to 8 workers live — visible speedup."),
    ]
    y0 = Inches(2.0)
    for i, (k, v) in enumerate(points):
        y = y0 + i * Inches(0.95)
        # number
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.65), y + Inches(0.05), Inches(0.5), Inches(0.5))
        nc.fill.solid(); nc.fill.fore_color.rgb = ACCENT
        nc.line.fill.background(); nc.shadow.inherit = False
        text(s, Inches(0.65), y + Inches(0.07), Inches(0.5), Inches(0.5),
             str(i+1), size=16, color=BG, bold=True, align=PP_ALIGN.CENTER)
        text(s, Inches(1.35), y, Inches(11.5), Inches(0.4),
             k, size=18, color=ACCENT, bold=True)
        text(s, Inches(1.35), y + Inches(0.4), Inches(11.5), Inches(0.5),
             v, size=14, color=FG)

# ---- main -----------------------------------------------------------------

slide_title()
slide_what()
slide_stack()
slide_components()
slide_pipeline()
slide_primitives()
slide_algorithm()
slide_balance()
slide_ui()
slide_run()
slide_summary()

out = "Matchmaking_Service_Presentation.pptx"
prs.save(out)
print(f"Wrote {out} ({len(prs.slides)} slides)")
